"""
增强文档加载器
支持更多格式、OCR、批量处理等功能
"""
import os
import logging
import tempfile
import shutil
from typing import List, Dict, Any, Optional, Tuple, Callable
from pathlib import Path
from dataclasses import dataclass
import hashlib

logger = logging.getLogger(__name__)


@dataclass
class LoadedDocument:
    """已加载的文档"""
    content: str
    metadata: Dict[str, Any]
    source: str
    
    def to_dict(self) -> Dict[str, Any]:
        return {
            "content": self.content,
            "metadata": self.metadata,
            "source": self.source,
        }


class EnhancedDocumentLoader:
    """
    增强文档加载器
    
    支持：
    1. 多格式解析（TXT, MD, DOCX, XLSX, PDF, CSV, JSON, XML, PPTX, 图片OCR）
    2. 批量处理
    3. 增量导入
    4. 文件变更检测
    5. 智能元数据提取
    """
    
    SUPPORTED_FORMATS = {
        # 文本格式
        ".txt", ".text", ".md", ".markdown", ".rst",
        # 文档格式
        ".docx", ".doc", ".pdf", ".rtf", ".odt",
        # 表格格式
        ".xlsx", ".xls", ".csv", ".ods",
        # 演示格式
        ".pptx", ".ppt", ".odp",
        # 数据格式
        ".json", ".xml", ".yaml", ".yml",
        # 图片（OCR）
        ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff",
    }
    
    def __init__(
        self,
        enable_ocr: bool = False,
        ocr_language: str = "chi_sim+eng",
        max_file_size: int = 10 * 1024 * 1024,  # 10MB
        encoding_detection: bool = True,
    ):
        """
        初始化增强文档加载器
        
        Args:
            enable_ocr: 是否启用 OCR
            ocr_language: OCR 语言
            max_file_size: 最大文件大小
            encoding_detection: 是否自动检测编码
        """
        self.enable_ocr = enable_ocr
        self.ocr_language = ocr_language
        self.max_file_size = max_file_size
        self.encoding_detection = encoding_detection
        
        # 文件状态追踪（用于增量导入）
        self._file_states: Dict[str, Dict[str, Any]] = {}
        
        logger.info(f"EnhancedDocumentLoader initialized (OCR: {enable_ocr})")
    
    def load_file(
        self,
        file_path: str,
        encoding: Optional[str] = None,
    ) -> Tuple[str, Dict[str, Any]]:
        """
        加载单个文件
        
        Args:
            file_path: 文件路径
            encoding: 文件编码（可选）
            
        Returns:
            (内容, 元数据) 元组
        """
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # 检查文件大小
        file_size = path.stat().st_size
        if file_size > self.max_file_size:
            raise ValueError(f"File too large: {file_size} > {self.max_file_size}")
        
        # 获取文件类型
        file_type = path.suffix.lower()
        
        # 根据类型选择加载方法
        if file_type in [".txt", ".text", ".md", ".markdown", ".rst"]:
            content = self._load_text_file(path, encoding)
        elif file_type == ".docx":
            content = self._load_docx(path)
        elif file_type == ".pdf":
            content = self._load_pdf(path)
        elif file_type in [".xlsx", ".xls"]:
            content = self._load_excel(path)
        elif file_type == ".csv":
            content = self._load_csv(path, encoding)
        elif file_type == ".json":
            content = self._load_json(path, encoding)
        elif file_type == ".xml":
            content = self._load_xml(path, encoding)
        elif file_type == ".pptx":
            content = self._load_pptx(path)
        elif file_type in [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"]:
            content = self._load_image_ocr(path)
        else:
            # 尝试作为文本加载
            content = self._load_text_file(path, encoding)
        
        # 提取元数据
        metadata = self._extract_metadata(path, file_type, file_size)
        
        # 更新文件状态
        self._update_file_state(path, content)
        
        return content, metadata
    
    def load_directory(
        self,
        directory: str,
        recursive: bool = True,
        incremental: bool = False,
        file_filter: Optional[Callable[[str], bool]] = None,
    ) -> List[Tuple[str, Dict[str, Any]]]:
        """
        加载目录下的所有文件
        
        Args:
            directory: 目录路径
            recursive: 是否递归
            incremental: 是否增量导入（只处理新文件）
            file_filter: 文件过滤器
            
        Returns:
            (内容, 元数据) 元组列表
        """
        dir_path = Path(directory)
        
        if not dir_path.exists():
            raise FileNotFoundError(f"Directory not found: {directory}")
        
        results = []
        
        # 遍历文件
        if recursive:
            files = dir_path.rglob("*")
        else:
            files = dir_path.glob("*")
        
        for file_path in files:
            # 跳过目录
            if not file_path.is_file():
                continue
            
            # 检查格式
            if file_path.suffix.lower() not in self.SUPPORTED_FORMATS:
                continue
            
            # 应用过滤器
            if file_filter and not file_filter(str(file_path)):
                continue
            
            # 增量导入检查
            if incremental and not self._is_file_changed(file_path):
                logger.debug(f"Skipping unchanged file: {file_path}")
                continue
            
            # 加载文件
            try:
                content, metadata = self.load_file(str(file_path))
                results.append((content, metadata))
                logger.info(f"Loaded: {file_path}")
            except Exception as e:
                logger.error(f"Failed to load {file_path}: {e}")
        
        logger.info(f"Loaded {len(results)} files from {directory}")
        return results
    
    def _load_text_file(
        self,
        path: Path,
        encoding: Optional[str] = None,
    ) -> str:
        """加载文本文件"""
        # 自动检测编码
        if encoding is None and self.encoding_detection:
            encoding = self._detect_encoding(path)
        
        encoding = encoding or "utf-8"
        
        try:
            with open(path, "r", encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            # 尝试其他编码
            for enc in ["gbk", "gb2312", "latin-1"]:
                try:
                    with open(path, "r", encoding=enc) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            raise ValueError(f"Cannot decode file with any encoding: {path}")
    
    def _load_docx(self, path: Path) -> str:
        """加载 Word 文档"""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("Please install python-docx: pip install python-docx")
        
        doc = Document(str(path))
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        return "\n\n".join(paragraphs)
    
    def _load_pdf(self, path: Path) -> str:
        """加载 PDF 文档"""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError("Please install pymupdf: pip install pymupdf")
        
        doc = fitz.open(str(path))
        pages = []
        for page_num, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                pages.append(f"[Page {page_num + 1}]\n{text}")
        doc.close()
        return "\n\n".join(pages)
    
    def _load_excel(self, path: Path) -> str:
        """加载 Excel 文件"""
        try:
            import pandas as pd
        except ImportError:
            raise ImportError("Please install pandas: pip install pandas openpyxl")
        
        df_dict = pd.read_excel(str(path), sheet_name=None)
        contents = []
        for sheet_name, df in df_dict.items():
            contents.append(f"## Sheet: {sheet_name}\n")
            contents.append(df.to_string(index=False))
            contents.append("\n")
        return "\n".join(contents)
    
    def _load_csv(
        self,
        path: Path,
        encoding: Optional[str] = None,
    ) -> str:
        """加载 CSV 文件"""
        try:
            import pandas as pd
        except ImportError:
            raise ImportError("Please install pandas: pip install pandas")
        
        encoding = encoding or "utf-8"
        df = pd.read_csv(str(path), encoding=encoding)
        return df.to_string(index=False)
    
    def _load_json(
        self,
        path: Path,
        encoding: Optional[str] = None,
    ) -> str:
        """加载 JSON 文件"""
        import json
        
        encoding = encoding or "utf-8"
        with open(path, "r", encoding=encoding) as f:
            data = json.load(f)
        
        # 格式化为可读文本
        return json.dumps(data, indent=2, ensure_ascii=False)
    
    def _load_xml(
        self,
        path: Path,
        encoding: Optional[str] = None,
    ) -> str:
        """加载 XML 文件"""
        try:
            import xmltodict
        except ImportError:
            raise ImportError("Please install xmltodict: pip install xmltodict")
        
        import json
        
        encoding = encoding or "utf-8"
        with open(path, "r", encoding=encoding) as f:
            data = xmltodict.parse(f.read())
        
        return json.dumps(data, indent=2, ensure_ascii=False)
    
    def _load_pptx(self, path: Path) -> str:
        """加载 PowerPoint 文档"""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("Please install python-pptx: pip install python-pptx")
        
        prs = Presentation(str(path))
        slides = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = [f"## Slide {slide_num}"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text)
            slides.append("\n".join(slide_content))
        
        return "\n\n".join(slides)
    
    def _load_image_ocr(self, path: Path) -> str:
        """加载图片并执行 OCR"""
        if not self.enable_ocr:
            raise ValueError("OCR is disabled")
        
        try:
            import pytesseract
            from PIL import Image
        except ImportError:
            raise ImportError("Please install pytesseract and Pillow: pip install pytesseract Pillow")
        
        img = Image.open(str(path))
        text = pytesseract.image_to_string(img, lang=self.ocr_language)
        return text
    
    def _detect_encoding(self, path: Path) -> str:
        """检测文件编码"""
        try:
            import chardet
        except ImportError:
            return "utf-8"
        
        with open(path, "rb") as f:
            raw = f.read(10000)  # 读取前 10KB
            result = chardet.detect(raw)
            return result.get("encoding", "utf-8")
    
    def _extract_metadata(
        self,
        path: Path,
        file_type: str,
        file_size: int,
    ) -> Dict[str, Any]:
        """提取文件元数据"""
        stat = path.stat()
        
        return {
            "filename": path.name,
            "filepath": str(path),
            "file_type": file_type,
            "file_size": file_size,
            "created_time": stat.st_ctime,
            "modified_time": stat.st_mtime,
            "format": file_type.lstrip("."),
        }
    
    def _compute_file_hash(self, path: Path) -> str:
        """计算文件哈希"""
        hasher = hashlib.md5()
        with open(path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hasher.update(chunk)
        return hasher.hexdigest()
    
    def _update_file_state(self, path: Path, content: str):
        """更新文件状态"""
        self._file_states[str(path)] = {
            "hash": self._compute_file_hash(path),
            "content_hash": hashlib.md5(content.encode()).hexdigest(),
            "size": path.stat().st_size,
            "mtime": path.stat().st_mtime,
        }
    
    def _is_file_changed(self, path: Path) -> bool:
        """检查文件是否已更改"""
        path_str = str(path)
        
        if path_str not in self._file_states:
            return True
        
        current_hash = self._compute_file_hash(path)
        stored_hash = self._file_states[path_str].get("hash")
        
        return current_hash != stored_hash
    
    def get_supported_formats(self) -> set:
        """获取支持的格式列表"""
        return self.SUPPORTED_FORMATS.copy()


# 全局加载器实例
_enhanced_loader: Optional[EnhancedDocumentLoader] = None


def get_enhanced_loader(
    enable_ocr: bool = False,
    **kwargs
) -> EnhancedDocumentLoader:
    """
    获取增强文档加载器（单例）
    
    Args:
        enable_ocr: 是否启用 OCR
        **kwargs: 其他参数
        
    Returns:
        加载器实例
    """
    global _enhanced_loader
    
    if _enhanced_loader is None:
        _enhanced_loader = EnhancedDocumentLoader(
            enable_ocr=enable_ocr,
            **kwargs
        )
    
    return _enhanced_loader
